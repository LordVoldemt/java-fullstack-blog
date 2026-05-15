from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "presentations"
HTML_PATH = OUT_DIR / "llm-basics-overview-ppt.html"
PPTX_PATH = OUT_DIR / "llm-basics-overview.pptx"


SLIDES = [
    {
        "kicker": "LLM Basics",
        "title": "大模型基础入门",
        "subtitle": "从 Prompt、Embedding、RAG 到工具调用与 MCP，先把地图看清。",
        "type": "cover",
        "tags": ["Prompt", "RAG", "Embedding", "MCP", "Agent"],
    },
    {
        "kicker": "Learning Path",
        "title": "先学让模型答得稳，再学让模型做更多",
        "subtitle": "一条适合初学者的递进路线",
        "type": "timeline",
        "items": [
            ("01", "LLM + Prompt", "理解模型如何根据上下文生成回答"),
            ("02", "Embedding + 检索", "把文本放进语义空间，找到相关内容"),
            ("03", "RAG", "先找知识，再让模型基于上下文回答"),
            ("04", "Function Calling", "让模型表达结构化的工具调用意图"),
            ("05", "Agent + MCP", "把任务推进和工具生态标准化连接起来"),
        ],
    },
    {
        "kicker": "Core Difference",
        "title": "大模型不是确定性函数",
        "subtitle": "它更像一个强大的概率生成器",
        "type": "compare",
        "leftTitle": "传统程序",
        "rightTitle": "大模型",
        "left": ["规则明确", "输入输出关系可精确设计", "行为边界相对稳定"],
        "right": ["理解上下文意图", "在模糊任务中生成合理答案", "结果具有概率性和不确定性"],
        "takeaway": "Prompt、RAG 和工具调用，本质上都在约束和利用这个生成器。",
    },
    {
        "kicker": "Concept Map",
        "title": "这些概念不是同一层东西",
        "subtitle": "它们各自解决链路上的不同问题",
        "type": "matrix",
        "items": [
            ("LLM", "理解和生成语言", "核心大脑"),
            ("Prompt", "描述目标和边界", "任务说明书"),
            ("Embedding", "把文本变成语义向量", "语义坐标"),
            ("向量检索", "寻找相近内容", "语义搜索引擎"),
            ("RAG", "检索后再生成", "知识增强链路"),
            ("MCP", "标准化接入工具", "工具连接协议"),
        ],
    },
    {
        "kicker": "Prompt",
        "title": "Prompt 是对模型行为的任务建模",
        "subtitle": "不是简单换一种问法，而是明确目标、上下文和输出约束。",
        "type": "three",
        "items": [
            ("目标", "你希望模型完成什么任务"),
            ("边界", "模型应该基于哪些上下文工作"),
            ("格式", "答案应该以什么结构和风格输出"),
        ],
        "note": "Prompt 能改善质量，但不能替代外部知识和实时数据。",
    },
    {
        "kicker": "Embedding",
        "title": "Embedding 让系统能按语义找内容",
        "subtitle": "从关键词匹配，走向“意思接近”的检索。",
        "type": "embedding",
        "pairs": [
            ("Java 并发", "多线程编程"),
            ("缓存雪崩", "大量缓存同时失效"),
            ("知识库问答", "基于文档回答问题"),
        ],
        "note": "先找对上下文，再让模型生成，是很多 AI 应用稳定的关键。",
    },
    {
        "kicker": "RAG Loop",
        "title": "最小 RAG 闭环：初学者最值得先跑通",
        "subtitle": "以 Java 学习知识库问答为例",
        "type": "flow",
        "steps": [
            "准备文档",
            "切片",
            "生成 Embedding",
            "存入向量库",
            "检索 Top-K",
            "拼接 Prompt",
            "模型回答",
        ],
        "formula": "documents -> chunk -> embedding -> vector store -> retrieve -> prompt -> answer",
    },
    {
        "kicker": "Tools & MCP",
        "title": "工具调用把模型接到真实世界",
        "subtitle": "模型负责判断要做什么，外部工具负责真正执行。",
        "type": "tools",
        "items": [
            ("Function Calling", "让模型用结构化格式表达调用哪个工具、传什么参数"),
            ("工具调用", "查天气、读文件、查数据库、调业务接口、执行代码"),
            ("MCP", "用更统一的协议描述模型可发现、可使用的外部能力"),
        ],
        "note": "MCP 不是替代 RAG 或 Prompt，而是让工具连接更标准。",
    },
    {
        "kicker": "Takeaway",
        "title": "看清整张图，再走最小闭环",
        "subtitle": "不要一上来追复杂 Agent，先让检索、上下文和回答质量稳定。",
        "type": "summary",
        "items": [
            ("Prompt", "表达任务"),
            ("Embedding", "找到知识"),
            ("RAG", "增强回答"),
            ("Function Calling", "连接动作"),
            ("MCP", "标准接入"),
        ],
        "closing": "地图清楚了，AI 应用开发才不会每走一步都像在摸黑。",
    },
]


COLORS = {
    "bg": "07130F",
    "panel": "0D211A",
    "panel2": "102A22",
    "mint": "06D6A0",
    "cyan": "22D3EE",
    "amber": "F59E0B",
    "red": "FB7185",
    "text": "EAFBF5",
    "muted": "9CC7B8",
    "line": "1E4D40",
}


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def add_text(slide, text, x, y, w, h, size=18, color="text", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS.get(fill, fill))
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(COLORS.get(line or fill, line or fill))
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x, y, w, h, color="line", width=1.4):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(width)
    return line


def add_pill(slide, text, x, y, w, accent="mint"):
    add_rect(slide, x, y, w, 0.34, "panel2", accent, radius=True, transparency=8)
    add_text(slide, text, x + 0.12, y + 0.08, w - 0.24, 0.2, 9.5, accent, True, "center")


def decorate(slide, index):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
    add_rect(slide, 0, 0, 10, 0.08, "mint", "mint")
    add_rect(slide, 0, 5.49, 10, 0.135, "panel", "panel")
    add_text(slide, f"{index:02d} / {len(SLIDES):02d}", 8.95, 5.21, 0.72, 0.18, 8.5, "muted", False, "right")


def add_header(slide, data, index):
    decorate(slide, index)
    add_text(slide, data["kicker"].upper(), 0.55, 0.34, 2.3, 0.24, 9.5, "mint", True)
    add_text(slide, data["title"], 0.55, 0.68, 8.8, 0.55, 25, "text", True)
    add_text(slide, data["subtitle"], 0.57, 1.22, 8.6, 0.36, 12.5, "muted")


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    for i, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if data["type"] == "cover":
            decorate(slide, i)
            add_text(slide, data["kicker"], 0.65, 0.62, 2.6, 0.28, 12, "mint", True)
            add_text(slide, data["title"], 0.62, 1.13, 5.6, 0.78, 34, "text", True)
            add_text(slide, data["subtitle"], 0.67, 2.02, 5.4, 0.68, 16, "muted")
            for n, tag in enumerate(data["tags"]):
                add_pill(slide, tag, 0.68 + n * 1.05, 3.15, 0.86 if tag != "Embedding" else 1.1, "cyan" if n % 2 else "mint")
            add_rect(slide, 6.65, 0.88, 2.35, 2.35, "panel", "line", radius=True, transparency=4)
            add_text(slide, "LLM", 7.28, 1.37, 1.1, 0.36, 27, "mint", True, "center")
            for angle, label in enumerate(["Prompt", "RAG", "Tool", "MCP"]):
                x = [6.35, 8.58, 6.46, 8.45][angle]
                y = [0.73, 0.76, 3.02, 3.03][angle]
                add_rect(slide, x, y, 0.8, 0.46, "panel2", "cyan", radius=True, transparency=4)
                add_text(slide, label, x + 0.08, y + 0.14, 0.64, 0.16, 8.5, "text", True, "center")
                add_line(slide, x + 0.4, y + 0.45, 7.83 - (x + 0.4), 1.94 - (y + 0.45), "line", 1.1)
        else:
            add_header(slide, data, i)
            t = data["type"]
            if t == "timeline":
                start_x = 0.76
                add_line(slide, 1.0, 2.35, 8.0, 0, "line", 2.2)
                for n, (num, title, desc) in enumerate(data["items"]):
                    x = start_x + n * 1.75
                    add_rect(slide, x, 1.92, 0.58, 0.58, "panel2", "mint", radius=True)
                    add_text(slide, num, x + 0.13, 2.08, 0.32, 0.16, 10, "mint", True, "center")
                    add_text(slide, title, x - 0.1, 2.72, 1.35, 0.3, 13, "text", True, "center")
                    add_text(slide, desc, x - 0.16, 3.12, 1.5, 0.72, 9.5, "muted", False, "center")
            elif t == "compare":
                add_rect(slide, 0.72, 1.9, 3.92, 2.35, "panel", "line", radius=True, transparency=4)
                add_rect(slide, 5.36, 1.9, 3.92, 2.35, "panel", "line", radius=True, transparency=4)
                add_text(slide, data["leftTitle"], 1.0, 2.12, 1.5, 0.28, 16, "amber", True)
                add_text(slide, data["rightTitle"], 5.64, 2.12, 1.5, 0.28, 16, "mint", True)
                for n, text in enumerate(data["left"]):
                    add_text(slide, text, 1.0, 2.65 + n * 0.42, 3.0, 0.24, 13, "text")
                for n, text in enumerate(data["right"]):
                    add_text(slide, text, 5.64, 2.65 + n * 0.42, 3.0, 0.24, 13, "text")
                add_text(slide, data["takeaway"], 1.35, 4.68, 7.3, 0.28, 13, "cyan", True, "center")
            elif t == "matrix":
                for n, (name, solve, metaphor) in enumerate(data["items"]):
                    col = n % 3
                    row = n // 3
                    x = 0.7 + col * 3.08
                    y = 1.92 + row * 1.3
                    add_rect(slide, x, y, 2.6, 0.93, "panel", "line", radius=True, transparency=5)
                    add_text(slide, name, x + 0.18, y + 0.16, 1.4, 0.22, 15, "mint", True)
                    add_text(slide, solve, x + 0.18, y + 0.46, 2.08, 0.18, 9.8, "text")
                    add_text(slide, metaphor, x + 0.18, y + 0.68, 1.9, 0.16, 8.8, "muted")
            elif t == "three":
                for n, (title, desc) in enumerate(data["items"]):
                    x = 0.8 + n * 3.0
                    add_rect(slide, x, 2.08, 2.38, 1.35, "panel", "line", radius=True, transparency=4)
                    add_text(slide, f"0{n + 1}", x + 0.18, 2.28, 0.4, 0.24, 12, "cyan", True)
                    add_text(slide, title, x + 0.66, 2.22, 1.1, 0.28, 18, "mint", True)
                    add_text(slide, desc, x + 0.28, 2.84, 1.78, 0.36, 11.2, "text", False, "center")
                add_text(slide, data["note"], 1.15, 4.32, 7.7, 0.3, 14, "amber", True, "center")
            elif t == "embedding":
                add_text(slide, "语义空间", 4.15, 1.96, 1.7, 0.28, 18, "mint", True, "center")
                add_rect(slide, 4.55, 2.5, 0.84, 0.84, "panel2", "mint", radius=True)
                add_text(slide, "向量", 4.77, 2.78, 0.4, 0.16, 10, "text", True, "center")
                for n, (a, b) in enumerate(data["pairs"]):
                    y = 2.0 + n * 0.72
                    add_rect(slide, 0.85, y, 2.25, 0.44, "panel", "cyan", radius=True, transparency=6)
                    add_rect(slide, 6.88, y, 2.25, 0.44, "panel", "cyan", radius=True, transparency=6)
                    add_text(slide, a, 1.02, y + 0.13, 1.88, 0.16, 10.5, "text", True, "center")
                    add_text(slide, b, 7.05, y + 0.13, 1.88, 0.16, 10.5, "text", True, "center")
                    add_line(slide, 3.12, y + 0.22, 1.32, 0.48 - n * 0.05, "line", 1.1)
                    add_line(slide, 5.42, y + 0.68 - n * 0.05, 1.36, -0.46 + n * 0.05, "line", 1.1)
                add_text(slide, data["note"], 1.2, 4.55, 7.6, 0.3, 13, "cyan", True, "center")
            elif t == "flow":
                for n, step in enumerate(data["steps"]):
                    x = 0.48 + (n % 4) * 2.34
                    y = 1.92 + (n // 4) * 1.08
                    add_rect(slide, x, y, 1.75, 0.6, "panel", "line", radius=True, transparency=4)
                    add_text(slide, f"{n + 1}", x + 0.15, y + 0.2, 0.2, 0.16, 10, "mint", True)
                    add_text(slide, step, x + 0.42, y + 0.19, 1.1, 0.17, 10.5, "text", True)
                    if n < len(data["steps"]) - 1 and n != 3:
                        add_line(slide, x + 1.82, y + 0.3, 0.42, 0, "cyan", 1.5)
                add_rect(slide, 0.9, 4.22, 8.2, 0.5, "panel2", "line", radius=True, transparency=3)
                add_text(slide, data["formula"], 1.05, 4.41, 7.9, 0.14, 9.5, "cyan", False, "center")
            elif t == "tools":
                for n, (title, desc) in enumerate(data["items"]):
                    y = 1.9 + n * 0.78
                    add_text(slide, title, 0.9, y, 2.2, 0.24, 15, "mint", True)
                    add_line(slide, 3.15, y + 0.15, 0.72, 0, "cyan", 1.8)
                    add_text(slide, desc, 4.05, y, 4.85, 0.36, 11.5, "text")
                add_text(slide, data["note"], 1.2, 4.55, 7.6, 0.3, 13, "amber", True, "center")
            elif t == "summary":
                for n, (title, desc) in enumerate(data["items"]):
                    x = 0.72 + n * 1.84
                    add_rect(slide, x, 2.0, 1.42, 1.18, "panel", "line", radius=True, transparency=5)
                    add_text(slide, title, x + 0.12, 2.28, 1.18, 0.18, 11.5, "mint", True, "center")
                    add_text(slide, desc, x + 0.16, 2.73, 1.08, 0.16, 10, "text", False, "center")
                add_text(slide, data["closing"], 1.0, 4.2, 8.0, 0.36, 15, "cyan", True, "center")

    prs.save(PPTX_PATH)


def html_escape(text):
    return (
        str(text)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def render_slide_html(data, index):
    cls = f"slide slide-{data['type']}"
    h = [f'<section class="{cls}" data-index="{index}">']
    h.append('<div class="top-line"></div>')
    h.append(f'<div class="page-no">{index:02d} / {len(SLIDES):02d}</div>')
    if data["type"] == "cover":
        h.append('<div class="cover-grid anim-item">')
        h.append('<div class="cover-copy">')
        h.append(f'<p class="kicker an">{html_escape(data["kicker"])}</p>')
        h.append(f'<h1 class="an">{html_escape(data["title"])}</h1>')
        h.append(f'<p class="subtitle an">{html_escape(data["subtitle"])}</p>')
        h.append('<div class="tag-row an">')
        for tag in data["tags"]:
            h.append(f'<span>{html_escape(tag)}</span>')
        h.append('</div></div>')
        h.append('<div class="orbit an"><div class="core">LLM</div><span>Prompt</span><span>RAG</span><span>Tool</span><span>MCP</span></div>')
        h.append('</div>')
    else:
        h.append(f'<p class="kicker an">{html_escape(data["kicker"])}</p>')
        h.append(f'<h2 class="an">{html_escape(data["title"])}</h2>')
        h.append(f'<p class="subtitle an">{html_escape(data["subtitle"])}</p>')
        t = data["type"]
        if t == "timeline":
            h.append('<div class="timeline anim-item">')
            for num, title, desc in data["items"]:
                h.append(f'<div class="step"><b>{num}</b><strong>{html_escape(title)}</strong><p>{html_escape(desc)}</p></div>')
            h.append('</div>')
        elif t == "compare":
            h.append('<div class="compare anim-item">')
            h.append(f'<div><h3>{html_escape(data["leftTitle"])}</h3>')
            h.extend(f'<p>{html_escape(x)}</p>' for x in data["left"])
            h.append('</div>')
            h.append(f'<div><h3>{html_escape(data["rightTitle"])}</h3>')
            h.extend(f'<p>{html_escape(x)}</p>' for x in data["right"])
            h.append('</div></div>')
            h.append(f'<div class="takeaway anim-item">{html_escape(data["takeaway"])}</div>')
        elif t == "matrix":
            h.append('<div class="matrix anim-item">')
            for name, solve, metaphor in data["items"]:
                h.append(f'<div><strong>{html_escape(name)}</strong><p>{html_escape(solve)}</p><small>{html_escape(metaphor)}</small></div>')
            h.append('</div>')
        elif t == "three":
            h.append('<div class="three anim-item">')
            for i, (title, desc) in enumerate(data["items"], start=1):
                h.append(f'<div><span>0{i}</span><strong>{html_escape(title)}</strong><p>{html_escape(desc)}</p></div>')
            h.append('</div>')
            h.append(f'<div class="note anim-item">{html_escape(data["note"])}</div>')
        elif t == "embedding":
            h.append('<div class="embedding-map anim-item"><div class="vector-core">语义空间<br><b>Embedding</b></div>')
            for a, b in data["pairs"]:
                h.append(f'<div class="pair"><span>{html_escape(a)}</span><i></i><span>{html_escape(b)}</span></div>')
            h.append('</div>')
            h.append(f'<div class="note anim-item">{html_escape(data["note"])}</div>')
        elif t == "flow":
            h.append('<div class="flow anim-item">')
            for i, step in enumerate(data["steps"], start=1):
                h.append(f'<div><span>{i}</span>{html_escape(step)}</div>')
            h.append('</div>')
            h.append(f'<code class="formula anim-item">{html_escape(data["formula"])}</code>')
        elif t == "tools":
            h.append('<div class="tool-list anim-item">')
            for title, desc in data["items"]:
                h.append(f'<div><strong>{html_escape(title)}</strong><p>{html_escape(desc)}</p></div>')
            h.append('</div>')
            h.append(f'<div class="note anim-item">{html_escape(data["note"])}</div>')
        elif t == "summary":
            h.append('<div class="summary-list anim-item">')
            for title, desc in data["items"]:
                h.append(f'<div><strong>{html_escape(title)}</strong><span>{html_escape(desc)}</span></div>')
            h.append('</div>')
            h.append(f'<div class="closing anim-item">{html_escape(data["closing"])}</div>')
    h.append('</section>')
    return "\n".join(h)


def build_html():
    slides_html = "\n".join(render_slide_html(s, i) for i, s in enumerate(SLIDES, start=1))
    html = f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>大模型基础入门 PPT</title>
<style>
:root {{
  --bg: #07130f;
  --panel: #0d211a;
  --panel2: #102a22;
  --mint: #06d6a0;
  --cyan: #22d3ee;
  --amber: #f59e0b;
  --text: #eafbF5;
  --muted: #9cc7b8;
  --line: #1e4d40;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  background: #030806;
  color: var(--text);
  font-family: "Microsoft YaHei", "PingFang SC", "Noto Sans SC", Arial, sans-serif;
  overflow: hidden;
}}
.deck {{
  width: 100vw;
  height: 100vh;
  display: grid;
  place-items: center;
}}
.slide {{
  display: none;
  position: relative;
  width: min(100vw, 177.78vh);
  height: min(56.25vw, 100vh);
  padding: 5.4% 6.2%;
  overflow: hidden;
  background:
    radial-gradient(circle at 84% 18%, rgba(34, 211, 238, .17), transparent 30%),
    radial-gradient(circle at 10% 86%, rgba(6, 214, 160, .13), transparent 28%),
    linear-gradient(135deg, #07130f 0%, #0a1713 52%, #07110f 100%);
}}
.slide.active {{ display: block; }}
.top-line {{ position:absolute; inset:0 auto auto 0; height:8px; width:100%; background:linear-gradient(90deg,var(--mint),var(--cyan)); }}
.page-no {{ position:absolute; right:6%; bottom:5.2%; color:var(--muted); font-size:clamp(10px,1.1vw,15px); }}
.kicker {{ margin:0 0 .65rem; color:var(--mint); font-weight:800; font-size:clamp(13px,1.35vw,18px); text-transform:uppercase; }}
h1,h2 {{ margin:0; line-height:1.08; letter-spacing:0; }}
h1 {{ font-size:clamp(42px,6vw,86px); max-width:65%; }}
h2 {{ font-size:clamp(25px,3.5vw,50px); max-width:88%; }}
.subtitle {{ margin:.8rem 0 0; color:var(--muted); max-width:78%; font-size:clamp(14px,1.45vw,22px); line-height:1.5; }}
.cover-grid {{ display:grid; grid-template-columns: 1.35fr .9fr; gap:6%; height:100%; align-items:center; }}
.tag-row {{ display:flex; flex-wrap:wrap; gap:.65rem; margin-top:2rem; }}
.tag-row span, .slide .step b {{
  border:1px solid rgba(34,211,238,.52);
  background:rgba(16,42,34,.78);
  border-radius:999px;
  color:var(--cyan);
  padding:.38rem .7rem;
  font-weight:800;
}}
.orbit {{ position:relative; height:min(45vw,420px); border:1px solid rgba(34,211,238,.3); border-radius:24px; background:rgba(13,33,26,.76); }}
.core {{ position:absolute; inset:38% 30%; display:grid; place-items:center; color:var(--mint); font-size:clamp(28px,4vw,54px); font-weight:900; border:1px solid var(--line); border-radius:18px; }}
.orbit span {{ position:absolute; color:var(--text); background:var(--panel2); border:1px solid var(--cyan); border-radius:999px; padding:.5rem .75rem; font-weight:800; }}
.orbit span:nth-child(2) {{ left:8%; top:12%; }}
.orbit span:nth-child(3) {{ right:8%; top:14%; }}
.orbit span:nth-child(4) {{ left:10%; bottom:14%; }}
.orbit span:nth-child(5) {{ right:10%; bottom:13%; }}
.timeline {{ margin-top:8%; display:grid; grid-template-columns:repeat(5,1fr); gap:1.1rem; position:relative; }}
.timeline:before {{ content:""; position:absolute; left:5%; right:5%; top:2.1rem; border-top:2px solid var(--line); }}
.step {{ position:relative; text-align:center; }}
.step b {{ display:inline-grid; place-items:center; width:3.9rem; height:3.9rem; padding:0; color:var(--mint); }}
.step strong {{ display:block; margin:1rem 0 .45rem; font-size:clamp(14px,1.45vw,20px); }}
.step p, .matrix p, .matrix small, .tool-list p, .three p {{ color:var(--muted); line-height:1.42; font-size:clamp(11px,1.1vw,16px); }}
.compare {{ margin-top:5%; display:grid; grid-template-columns:1fr 1fr; gap:6%; }}
.compare div, .matrix div, .three div, .tool-list div, .summary-list div, .flow div {{
  background:rgba(13,33,26,.78);
  border:1px solid var(--line);
  border-radius:8px;
}}
.compare div {{ padding:1.7rem; min-height:15rem; }}
.compare h3 {{ color:var(--mint); margin:0 0 1rem; font-size:clamp(20px,2vw,30px); }}
.compare p {{ font-size:clamp(15px,1.45vw,21px); }}
.takeaway,.note,.closing {{
  margin:2rem auto 0;
  text-align:center;
  color:var(--cyan);
  font-weight:900;
  font-size:clamp(17px,1.7vw,25px);
}}
.matrix {{ margin-top:4.5%; display:grid; grid-template-columns:repeat(3,1fr); gap:1rem; }}
.matrix div {{ min-height:8rem; padding:1.2rem; }}
.matrix strong, .tool-list strong, .summary-list strong {{ color:var(--mint); font-size:clamp(18px,1.8vw,25px); }}
.three {{ display:grid; grid-template-columns:repeat(3,1fr); gap:1.2rem; margin-top:6%; }}
.three div {{ min-height:12rem; padding:1.5rem; text-align:center; }}
.three span {{ color:var(--cyan); font-weight:900; }}
.three strong {{ display:block; margin:.8rem 0 1rem; color:var(--mint); font-size:clamp(28px,2.8vw,44px); }}
.embedding-map {{ margin-top:4%; display:grid; gap:.75rem; grid-template-columns:1fr .8fr 1fr; align-items:center; }}
.vector-core {{ grid-row:1/4; grid-column:2; text-align:center; color:var(--muted); font-size:clamp(16px,1.6vw,24px); }}
.vector-core b {{ color:var(--mint); }}
.pair {{ display:contents; }}
.pair span {{ background:rgba(13,33,26,.85); border:1px solid var(--cyan); border-radius:999px; padding:.8rem; text-align:center; font-weight:800; }}
.pair i {{ min-height:2px; background:linear-gradient(90deg,var(--line),var(--cyan),var(--line)); }}
.flow {{ margin-top:4.8%; display:grid; grid-template-columns:repeat(4,1fr); gap:1rem; }}
.flow div {{ padding:1rem; min-height:4rem; font-weight:800; }}
.flow span {{ color:var(--mint); margin-right:.55rem; }}
.formula {{ display:block; margin:1.6rem auto 0; color:var(--cyan); background:rgba(16,42,34,.8); border:1px solid var(--line); border-radius:8px; padding:.9rem; text-align:center; max-width:88%; font-size:clamp(12px,1.25vw,18px); white-space:normal; }}
.tool-list {{ margin-top:4.5%; display:grid; gap:1rem; }}
.tool-list div {{ padding:1.1rem 1.4rem; display:grid; grid-template-columns:2.2fr 5fr; gap:1rem; align-items:center; }}
.summary-list {{ margin-top:5%; display:grid; grid-template-columns:repeat(5,1fr); gap:.85rem; }}
.summary-list div {{ min-height:8rem; padding:1.2rem .8rem; text-align:center; display:flex; flex-direction:column; justify-content:center; gap:.8rem; }}
.summary-list span {{ color:var(--muted); }}
.controls {{
  position:fixed;
  left:50%;
  bottom:1.1rem;
  transform:translateX(-50%);
  display:flex;
  gap:.55rem;
  z-index:5;
}}
button {{
  width:2.5rem;
  height:2.5rem;
  border-radius:999px;
  border:1px solid rgba(34,211,238,.35);
  background:rgba(7,19,15,.75);
  color:var(--text);
  font-weight:900;
  cursor:pointer;
}}
.an,.anim-item {{ opacity:0; transform:translateY(16px); animation:fadeUp .55s ease forwards; }}
.active .an:nth-of-type(2), .active .anim-item:nth-of-type(2) {{ animation-delay:.09s; }}
.active .an:nth-of-type(3), .active .anim-item:nth-of-type(3) {{ animation-delay:.16s; }}
@keyframes fadeUp {{ to {{ opacity:1; transform:translateY(0); }} }}
.static-export .an,
.static-export .anim-item {{
  opacity: 1 !important;
  transform: none !important;
  animation: none !important;
}}
@media (max-width: 760px) {{
  body {{ overflow:auto; }}
  .deck {{ min-height:100vh; height:auto; }}
  .slide {{ width:100vw; height:auto; min-height:100vh; padding:4.2rem 1.25rem 5rem; }}
  h1,h2,.subtitle {{ max-width:100%; }}
  .cover-grid,.timeline,.compare,.matrix,.three,.flow,.summary-list,.tool-list div {{ grid-template-columns:1fr; }}
  .timeline:before {{ display:none; }}
  .orbit {{ min-height:290px; }}
  .embedding-map {{ grid-template-columns:1fr; }}
  .vector-core {{ grid-row:auto; grid-column:auto; }}
  .pair {{ display:grid; grid-template-columns:1fr; gap:.5rem; }}
  .pair i {{ height:2px; }}
}}
</style>
</head>
<body>
<main class="deck">
{slides_html}
</main>
<nav class="controls" aria-label="幻灯片控制">
  <button type="button" id="prev" title="上一页">‹</button>
  <button type="button" id="next" title="下一页">›</button>
</nav>
<script>
const slides = Array.from(document.querySelectorAll('.slide'));
let current = 0;
function resetAnimations(slide) {{
  slide.querySelectorAll('.an, .anim-item, [style*="animation"]').forEach(function(item) {{
    const clone = item.cloneNode(true);
    item.parentNode.replaceChild(clone, item);
  }});
}}
function showSlide(index) {{
  current = (index + slides.length) % slides.length;
  slides.forEach((s, i) => s.classList.toggle('active', i === current));
  resetAnimations(slides[current]);
}}
document.getElementById('prev').addEventListener('click', () => showSlide(current - 1));
document.getElementById('next').addEventListener('click', () => showSlide(current + 1));
document.addEventListener('keydown', (event) => {{
  if (event.key === 'ArrowRight' || event.key === ' ') showSlide(current + 1);
  if (event.key === 'ArrowLeft') showSlide(current - 1);
}});
const hashSlide = Number.parseInt(location.hash.replace('#', ''), 10);
if (new URLSearchParams(location.search).has('static')) {{
  document.documentElement.classList.add('static-export');
}}
showSlide(Number.isFinite(hashSlide) ? hashSlide - 1 : 0);
</script>
</body>
</html>
"""
    HTML_PATH.write_text(html, encoding="utf-8")


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_html()
    build_pptx()
    print(HTML_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
